# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Blue & White Theme ──
BG = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF4, 0xF7, 0xFB)
BG_CARD = RGBColor(0xEE, 0xF2, 0xF9)
BLUE_DARK = RGBColor(0x1A, 0x36, 0x5D)
BLUE = RGBColor(0x1E, 0x56, 0xA0)
BLUE_MED = RGBColor(0x2E, 0x7D, 0xC9)
BLUE_LIGHT = RGBColor(0x4A, 0x9E, 0xE0)
BLUE_PALE = RGBColor(0xD6, 0xE6, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK = RGBColor(0x2D, 0x3A, 0x4A)
GRAY = RGBColor(0x5A, 0x6A, 0x7A)
LIGHT_GRAY = RGBColor(0x8A, 0x95, 0xA5)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
PHOTO = r"C:\Users\Sri\Pictures\Me\me.jpg"


def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def shape(slide, l, t, w, h, fill=BG_CARD, border=None, r=0.03):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1.2)
    else:
        s.line.fill.background()
    s.adjustments[0] = r
    return s


def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=18, col=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = col
    p.font.bold = bold; p.font.name = "Calibri"; p.alignment = align
    return tb


def bul(slide, l, t, w, h, items, sz=14, col=DARK, bc=BLUE):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5); p.space_before = Pt(2)
        r1 = p.add_run(); r1.text = "\u25b9 "; r1.font.size = Pt(sz); r1.font.color.rgb = bc; r1.font.name = "Calibri"
        r2 = p.add_run(); r2.text = item; r2.font.size = Pt(sz); r2.font.color.rgb = col; r2.font.name = "Calibri"
    return tb


def blue_line(slide, l, t, w=Inches(1.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()


def header(slide, title, sub=""):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)
    txt(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), title, 30, BLUE_DARK, True)
    if sub:
        txt(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4), sub, 15, BLUE_MED)
    blue_line(slide, Inches(0.8), Inches(1.3))


def section(title, sub=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_LIGHT)
    rect(s, Inches(0), Inches(0), Inches(0.15), Inches(7.5), BLUE)
    blue_line(s, Inches(4.5), Inches(3.2), Inches(4.3))
    txt(s, Inches(1), Inches(3.4), Inches(11.3), Inches(1), title, 40, BLUE_DARK, True, PP_ALIGN.CENTER)
    if sub:
        txt(s, Inches(1), Inches(4.5), Inches(11.3), Inches(0.6), sub, 17, GRAY, False, PP_ALIGN.CENTER)


def stat(slide, l, t, val, label, accent=BLUE):
    shape(slide, l, t, Inches(2.4), Inches(1.3), WHITE, accent, 0.05)
    txt(slide, l + Inches(0.1), t + Inches(0.1), Inches(2.2), Inches(0.55), val, 30, accent, True, PP_ALIGN.CENTER)
    txt(slide, l + Inches(0.1), t + Inches(0.75), Inches(2.2), Inches(0.4), label, 12, GRAY, False, PP_ALIGN.CENTER)


def arrow_flow(slide, l, t, w, steps, colors=None):
    """Draw a horizontal flow: box -> arrow -> box -> arrow -> box"""
    n = len(steps)
    box_w = (w - Inches(0.4) * (n - 1)) / n
    arrow_w = Inches(0.4)
    for i, step in enumerate(steps):
        x = l + (box_w + arrow_w) * i
        c = (colors[i] if colors else BLUE) if colors and i < len(colors) else BLUE
        shape(slide, x, t, box_w, Inches(0.65), c, None, 0.04)
        txt(slide, x + Inches(0.08), t + Inches(0.08), box_w - Inches(0.16), Inches(0.5), step, 10, WHITE, True, PP_ALIGN.CENTER)
        if i < n - 1:
            ax = x + box_w + Inches(0.05)
            txt(slide, ax, t + Inches(0.05), Inches(0.3), Inches(0.5), "\u2192", 22, BLUE, True, PP_ALIGN.CENTER)


def flow_vertical(slide, l, t, w, h, title, steps):
    """Vertical flow inside a card"""
    shape(slide, l, t, w, h, WHITE, BLUE, 0.03)
    txt(slide, l + Inches(0.2), t + Inches(0.1), w - Inches(0.4), Inches(0.35), title, 15, BLUE_DARK, True)
    step_h = (h - Inches(0.8)) / len(steps)
    for i, step in enumerate(steps):
        sy = t + Inches(0.55) + step_h * i
        shape(slide, l + Inches(0.2), sy, w - Inches(0.4), step_h - Inches(0.08), BG_CARD, BLUE_PALE, 0.02)
        txt(slide, l + Inches(0.35), sy + Inches(0.04), w - Inches(0.7), step_h - Inches(0.12), step, 11, DARK)
        if i < len(steps) - 1:
            txt(slide, l + w / 2 - Inches(0.15), sy + step_h - Inches(0.18), Inches(0.3), Inches(0.2), "\u2193", 14, BLUE_MED, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
rect(s, Inches(0), Inches(0), Inches(5), Inches(7.5), BLUE_DARK)
try:
    s.shapes.add_picture(PHOTO, Inches(1.5), Inches(1.0), Inches(2.2), Inches(2.2))
except:
    pass
txt(s, Inches(0.5), Inches(3.5), Inches(4), Inches(0.5), "Sri Hari Sirisipalli", 28, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(0.5), Inches(4.1), Inches(4), Inches(0.5), "AI Systems Engineer", 18, BLUE_LIGHT, False, PP_ALIGN.CENTER)
txt(s, Inches(0.5), Inches(4.7), Inches(4), Inches(0.3), "Visakhapatnam, India", 13, RGBColor(0xA0, 0xB8, 0xD0), False, PP_ALIGN.CENTER)
txt(s, Inches(0.5), Inches(5.1), Inches(4), Inches(0.3), "sriharisirisipalli0@gmail.com", 12, RGBColor(0xA0, 0xB8, 0xD0), False, PP_ALIGN.CENTER)

txt(s, Inches(5.8), Inches(1.5), Inches(6.5), Inches(1), "Scalable ML &\nRetrieval Architectures", 40, BLUE_DARK, True)
blue_line(s, Inches(5.8), Inches(3.0), Inches(2))
txt(s, Inches(5.8), Inches(3.3), Inches(6.8), Inches(1.2),
    "I build end-to-end AI systems \u2014 from raw data ingestion through production deployment. "
    "Specializing in retrieval architectures, digital twin modeling, and signal-driven anomaly detection "
    "on SOC2-compliant cloud infrastructure.",
    15, DARK)

stats_data = [("3+", "Years"), ("10+", "Projects"), ("6+", "Companies"), ("5", "Domains")]
for i, (v, l) in enumerate(stats_data):
    x = Inches(5.8) + Inches(i * 1.7)
    shape(s, x, Inches(5.0), Inches(1.5), Inches(1.1), BG_CARD, BLUE, 0.06)
    txt(s, x, Inches(5.05), Inches(1.5), Inches(0.55), v, 26, BLUE, True, PP_ALIGN.CENTER)
    txt(s, x, Inches(5.6), Inches(1.5), Inches(0.35), l, 11, GRAY, False, PP_ALIGN.CENTER)

txt(s, Inches(5.8), Inches(6.5), Inches(7), Inches(0.3),
    "srihari-sirisipalli.github.io  |  linkedin.com/in/sri-hari-sirisipalli  |  github.com/srihari-sirisipalli", 11, LIGHT_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 2: CAREER TIMELINE (REVERSE)
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "Career Timeline")

timeline = [
    ("Mar 2024 \u2013 Present", "Software Engineer \u2013 Pangeon (IntelliPat)", "LLM & ML Infra  \u2022  SOC2 AWS  \u2022  Patent AI  \u2022  Terraform IaC", BG_CARD),
    ("Jun 2023 \u2013 Present", "AI Systems Engineer \u2013 Independent Consulting & R&D", "Digital Twins  \u2022  Anomaly Detection  \u2022  CV/Retrieval  \u2022  RAG  \u2022  Agritech Advisory", BG_CARD),
    ("Mar \u2013 Jun 2023", "Data Engineer \u2013 Sas2Py", "SAS \u2192 PySpark migration  \u2022  >99% validation  \u2022  Dependency graph modeling", BLUE_PALE),
    ("Jul \u2013 Dec 2022", "ML Intern \u2013 Corteva Agriscience", "TF/PyTorch \u2192 ONNX  \u2022  Pruning & quantization  \u2022  Cross-platform deployment", BLUE_PALE),
    ("Jun \u2013 Dec 2022", "Software Eng Intern \u2013 Dojima Networks", "Polkadot cross-chain  \u2022  Prometheus\u2013Grafana monitoring  \u2022  Blockchain infra", BLUE_PALE),
]

for i, (period, role, desc, bg) in enumerate(timeline):
    y = Inches(1.6) + Inches(i * 1.1)
    rect(s, Inches(2.95), y + Inches(0.15), Pt(2), Inches(0.8 if i < len(timeline) - 1 else 0), BLUE_LIGHT)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.82), y + Inches(0.08), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = BLUE; dot.line.fill.background()
    txt(s, Inches(0.5), y, Inches(2.2), Inches(0.4), period, 13, BLUE, True, PP_ALIGN.RIGHT)
    shape(s, Inches(3.3), y, Inches(9.5), Inches(0.9), bg, BLUE_PALE, 0.02)
    txt(s, Inches(3.5), y + Inches(0.05), Inches(9), Inches(0.35), role, 14, BLUE_DARK, True)
    txt(s, Inches(3.5), y + Inches(0.42), Inches(9), Inches(0.4), desc, 11, GRAY)


# ════════════════════════════════════════════════════
# SECTION: WORK EXPERIENCE
# ════════════════════════════════════════════════════
section("Work Experience", "Professional roles at product companies and startups")

# ════════════════════════════════════════════════════
# SLIDE: PANGEON - Problem → Approach → Result
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "Pangeon (IntelliPat, Inc.)", "Software Engineer \u2013 LLM & ML Infrastructure  |  Mar 2024 \u2013 Present  |  Contract")

# Problem-Approach-Result cards
shape(s, Inches(0.6), Inches(1.5), Inches(3.7), Inches(2.8), WHITE, RGBColor(0xE7, 0x4C, 0x3C), 0.03)
txt(s, Inches(0.9), Inches(1.6), Inches(3.2), Inches(0.35), "Problem", 16, RGBColor(0xE7, 0x4C, 0x3C), True)
bul(s, Inches(0.9), Inches(2.0), Inches(3.2), Inches(2.0), [
    "Manual patent prior-art search is slow, expensive, and inconsistent",
    "EC2-heavy architecture \u2192 high costs, slow batch processing",
    "No reproducible infrastructure across environments",
], 12, DARK, RGBColor(0xE7, 0x4C, 0x3C))

shape(s, Inches(4.6), Inches(1.5), Inches(3.7), Inches(2.8), WHITE, BLUE, 0.03)
txt(s, Inches(4.9), Inches(1.6), Inches(3.2), Inches(0.35), "Approach", 16, BLUE, True)
bul(s, Inches(4.9), Inches(2.0), Inches(3.2), Inches(2.0), [
    "LLM semantic scoring + embedding pipelines for automated prior-art ranking",
    "Migrated compute to GPU-enabled Lambda with async processing",
    "Terraform IaC with SOC2-compliant security controls",
], 12, DARK, BLUE)

shape(s, Inches(8.6), Inches(1.5), Inches(4.1), Inches(2.8), WHITE, GREEN, 0.03)
txt(s, Inches(8.9), Inches(1.6), Inches(3.5), Inches(0.35), "Result / ROI", 16, GREEN, True)
bul(s, Inches(8.9), Inches(2.0), Inches(3.5), Inches(2.0), [
    "30% infrastructure cost reduction",
    "5\u20137x faster batch completion",
    "SOC2 Type I audit-ready architecture",
    "Reproducible cross-account deployments",
], 12, DARK, GREEN)

# Flow
txt(s, Inches(0.8), Inches(4.6), Inches(5), Inches(0.35), "System Architecture Flow", 15, BLUE_DARK, True)
arrow_flow(s, Inches(0.8), Inches(5.0), Inches(11.5), [
    "Patent\nIntake", "Embedding\nPipeline", "LLM Semantic\nScoring", "FAISS Vector\nSearch", "Prior-Art\nRanking", "Novelty\nReport"
], [BLUE_DARK, BLUE, BLUE_MED, BLUE_LIGHT, BLUE_MED, BLUE_DARK])

# Tech
shape(s, Inches(0.6), Inches(5.9), Inches(12), Inches(0.7), BG_CARD, BLUE_PALE, 0.02)
txt(s, Inches(0.9), Inches(5.95), Inches(1.2), Inches(0.3), "Stack", 12, BLUE, True)
txt(s, Inches(2.1), Inches(5.95), Inches(10), Inches(0.55),
    "Python  \u2022  LLMs  \u2022  FAISS  \u2022  AWS (EC2, Lambda, SQS, ECS, Amplify, GuardDuty, Secrets Mgr)  \u2022  Terraform  \u2022  FastAPI  \u2022  Vanta  \u2022  Docker",
    11, GRAY)


# ════════════════════════════════════════════════════
# SLIDE: SAS2PY
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "Sas2Py", "Data Engineer  |  Mar 2023 \u2013 Jun 2023  |  Full-Time")

# Problem-Approach-Result
shape(s, Inches(0.6), Inches(1.5), Inches(3.7), Inches(2.0), WHITE, RGBColor(0xE7, 0x4C, 0x3C), 0.03)
txt(s, Inches(0.9), Inches(1.6), Inches(3.2), Inches(0.35), "Problem", 16, RGBColor(0xE7, 0x4C, 0x3C), True)
bul(s, Inches(0.9), Inches(2.0), Inches(3.2), Inches(1.2), [
    "Legacy SAS pipelines \u2014 not scalable, single-threaded",
    "No automated validation between old and new outputs",
], 12, DARK, RGBColor(0xE7, 0x4C, 0x3C))

shape(s, Inches(4.6), Inches(1.5), Inches(3.7), Inches(2.0), WHITE, BLUE, 0.03)
txt(s, Inches(4.9), Inches(1.6), Inches(3.2), Inches(0.35), "Approach", 16, BLUE, True)
bul(s, Inches(4.9), Inches(2.0), Inches(3.2), Inches(1.2), [
    "Redesigned as PySpark distributed batch workflows",
    "Graph theory for dependency mapping & bottleneck detection",
], 12, DARK, BLUE)

shape(s, Inches(8.6), Inches(1.5), Inches(4.1), Inches(2.0), WHITE, GREEN, 0.03)
txt(s, Inches(8.9), Inches(1.6), Inches(3.5), Inches(0.35), "Result", 16, GREEN, True)
bul(s, Inches(8.9), Inches(2.0), Inches(3.5), Inches(1.2), [
    ">99% functional parity validated",
    "Eliminated scheduling bottlenecks",
    "Improved scalability & throughput",
], 12, DARK, GREEN)

# Flow
txt(s, Inches(0.8), Inches(3.8), Inches(5), Inches(0.35), "Migration Pipeline", 15, BLUE_DARK, True)
arrow_flow(s, Inches(0.8), Inches(4.2), Inches(11.5), [
    "Legacy SAS\nPipelines", "Dependency\nGraph Analysis", "PySpark\nRedesign", "Automated\nValidation", "Production\nDeployment"
], [BLUE_DARK, BLUE, BLUE_MED, BLUE_LIGHT, GREEN])


# ════════════════════════════════════════════════════
# SLIDE: INTERNSHIPS
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "Internships")

# Corteva
shape(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(4.0), WHITE, BLUE, 0.03)
txt(s, Inches(0.9), Inches(1.6), Inches(5.2), Inches(0.35), "Machine Learning Intern \u2013 Corteva Agriscience", 16, BLUE_DARK, True)
txt(s, Inches(0.9), Inches(2.0), Inches(5.2), Inches(0.25), "Jul 2022 \u2013 Dec 2022  |  Hyderabad", 12, BLUE_MED)

txt(s, Inches(0.9), Inches(2.4), Inches(2.4), Inches(0.3), "Problem", 12, RGBColor(0xE7, 0x4C, 0x3C), True)
txt(s, Inches(0.9), Inches(2.7), Inches(5.2), Inches(0.5), "Models locked to single framework, slow inference, no cross-platform deployment path", 11, DARK)

txt(s, Inches(0.9), Inches(3.2), Inches(2.4), Inches(0.3), "Approach \u2192 Result", 12, GREEN, True)
bul(s, Inches(0.9), Inches(3.5), Inches(5.2), Inches(1.5), [
    "TF/PyTorch \u2192 ONNX conversion for cross-platform compatibility",
    "Structured pruning + post-training quantization for inference speed",
    "Numerical consistency validation across all model formats",
], 12, DARK, GREEN)

# Dojima
shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(4.0), WHITE, BLUE, 0.03)
txt(s, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.35), "Software Engineer Intern \u2013 Dojima Networks", 16, BLUE_DARK, True)
txt(s, Inches(7.1), Inches(2.0), Inches(5.2), Inches(0.25), "Jun 2022 \u2013 Dec 2022  |  Remote", 12, BLUE_MED)

txt(s, Inches(7.1), Inches(2.4), Inches(2.4), Inches(0.3), "Problem", 12, RGBColor(0xE7, 0x4C, 0x3C), True)
txt(s, Inches(7.1), Inches(2.7), Inches(5.2), Inches(0.5), "No cross-chain interoperability, no real-time observability into API health", 11, DARK)

txt(s, Inches(7.1), Inches(3.2), Inches(2.4), Inches(0.3), "Approach \u2192 Result", 12, GREEN, True)
bul(s, Inches(7.1), Inches(3.5), Inches(5.2), Inches(1.5), [
    "Polkadot ecosystem integration for cross-chain interoperability",
    "Cross-chain transaction routing & state synchronization",
    "Prometheus\u2013Grafana stack for real-time API monitoring",
], 12, DARK, GREEN)


# ════════════════════════════════════════════════════
# SECTION: CONSULTING
# ════════════════════════════════════════════════════
section("Independent Consulting & R&D", "AI Systems Engineer  \u2022  Jun 2023 \u2013 Present  \u2022  Contract")

# ════════════════════════════════════════════════════
# SLIDE: DIGITAL TWIN
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "Digital Twin Systems \u2013 Offshore Riser Modeling", "Problem \u2192 Approach \u2192 Result")

# Problem
shape(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.0), WHITE, RGBColor(0xE7, 0x4C, 0x3C), 0.02)
txt(s, Inches(0.9), Inches(1.55), Inches(1.5), Inches(0.3), "Problem", 14, RGBColor(0xE7, 0x4C, 0x3C), True)
txt(s, Inches(2.5), Inches(1.55), Inches(9.8), Inches(0.8),
    "Running full physics simulations (OrcaFlex) for each new sea state takes hours per scenario. Engineers needed predictions across "
    "88,560 sea state combinations for fatigue analysis \u2014 infeasible with brute-force simulation. Manual feature selection was suboptimal.",
    12, DARK)

# Approach - Flow
txt(s, Inches(0.8), Inches(2.7), Inches(5), Inches(0.35), "Approach", 14, BLUE, True)
arrow_flow(s, Inches(0.8), Inches(3.1), Inches(11.5), [
    "88K Sea State\nSimulations", "3,888 Time-Series\nFiles (4Hz)", "156 Feature\nEngineering", "3,525 ML/DL\nExperiments", "~1.08M\nModel Fits", "Production\nPrediction"
], [BLUE_DARK, BLUE, BLUE_MED, BLUE_LIGHT, BLUE_MED, GREEN])

# Technical details
shape(s, Inches(0.6), Inches(4.0), Inches(5.8), Inches(2.5), BG_CARD, BLUE_PALE, 0.02)
txt(s, Inches(0.9), Inches(4.1), Inches(5.2), Inches(0.3), "Technical Approach", 14, BLUE_DARK, True)
bul(s, Inches(0.9), Inches(4.4), Inches(5.2), Inches(2.0), [
    "Surrogate ML models replacing physics simulation",
    "Circular regression for angular wave heading prediction",
    "Hierarchical classification-regression architecture",
    "8 scalers \u00d7 10 model types benchmarked systematically",
    "Strict dataset partitioning with leakage detection",
], 12, DARK)

# Results
shape(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.5), WHITE, GREEN, 0.02)
txt(s, Inches(7.1), Inches(4.1), Inches(5.2), Inches(0.3), "ROI & Results", 14, GREEN, True)
bul(s, Inches(7.1), Inches(4.4), Inches(5.2), Inches(2.0), [
    "R\u00b2 = 0.9992 on significant wave height (Random Forest)",
    "97.17% direction prediction accuracy",
    "41% worst-case angular error reduction",
    "Hours of simulation \u2192 seconds of prediction",
    "Fully reproducible retrain-to-predict pipeline",
], 12, DARK, GREEN)

# ROI callout
shape(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5), BLUE_DARK, None, 0.02)
txt(s, Inches(0.9), Inches(6.72), Inches(11.5), Inches(0.35),
    "ROI:  Replaced days of OrcaFlex simulation with sub-second ML predictions  \u2022  3,525 experiments automated end-to-end  \u2022  Production-ready pipeline",
    12, WHITE, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE: SIMULATION AUTOMATION
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "Engineering Simulation Automation", "ANSYS APDL + MOSES  \u2014  Automating Large-Scale Simulation Data Generation")

# MOSES Flow
txt(s, Inches(0.8), Inches(1.5), Inches(5), Inches(0.35), "MOSES Marine Simulation Pipeline", 15, BLUE_DARK, True)
flow_vertical(s, Inches(0.6), Inches(1.9), Inches(5.8), Inches(4.5), "", [
    "Excel parametric inputs: wave heights, spectra (PM, JONSWAP, ISSC, OCHI), headings, weight groups",
    "GL Noble Denton Tp derivation: \u221a(13\u00b7Hs) to \u221a(30\u00b7Hs)",
    "Template token replacement \u2192 MOSES .dat and .cif files",
    "Batch MOSES execution (Cartesian product of all parameters)",
    "6-DOF extraction: surge, sway, heave, roll, pitch, yaw",
    "Results aggregation + progress tracking + CSV summary",
])

# ANSYS Flow
txt(s, Inches(7.0), Inches(1.5), Inches(5), Inches(0.35), "ANSYS APDL \u2013 Mono Pile Turbine", 15, BLUE_DARK, True)
flow_vertical(s, Inches(6.8), Inches(1.9), Inches(5.8), Inches(4.5), "", [
    "Excel parametric inputs: geometry dimensions, loading conditions",
    "Python script generates APDL command files",
    "Automated mesh generation for mono pile structures",
    "Batch FEA structural analysis execution",
    "Output parsing for stress, displacement, modal results",
    "Systematic parametric studies across design configs",
])

shape(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.5), BLUE_DARK, None, 0.02)
txt(s, Inches(0.9), Inches(6.62), Inches(11.5), Inches(0.35),
    "ROI:  80%+ reduction in manual intervention  \u2022  Standardized, repeatable simulation workflows  \u2022  ML-ready dataset generation at scale",
    12, WHITE, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE: ANOMALY DETECTION
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "Anomaly Detection & Signal Processing", "Unsupervised Fault Detection for Naval Systems")

# Problem
shape(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.9), WHITE, RGBColor(0xE7, 0x4C, 0x3C), 0.02)
txt(s, Inches(0.9), Inches(1.55), Inches(1.5), Inches(0.3), "Problem", 14, RGBColor(0xE7, 0x4C, 0x3C), True)
txt(s, Inches(2.5), Inches(1.55), Inches(9.8), Inches(0.7),
    "Naval equipment failures are costly and dangerous. Labeled fault data is scarce. "
    "Need to detect anomalies from vibration sensor data without supervised training labels.",
    12, DARK)

# Flow
txt(s, Inches(0.8), Inches(2.6), Inches(5), Inches(0.3), "Detection Pipeline", 14, BLUE_DARK, True)
arrow_flow(s, Inches(0.8), Inches(2.95), Inches(11.5), [
    "Multi-Day\nSensor Data", "16+ Time\nFeatures", "11 Freq\nFeatures", "PCA\nReduction", "K-Means\nClustering", "Anomaly\nScoring"
], [BLUE_DARK, BLUE, BLUE, BLUE_MED, BLUE_LIGHT, GREEN])

# Technical + Results side by side
shape(s, Inches(0.6), Inches(3.8), Inches(5.8), Inches(2.8), BG_CARD, BLUE_PALE, 0.02)
txt(s, Inches(0.9), Inches(3.9), Inches(5.2), Inches(0.3), "Technical Approach", 14, BLUE_DARK, True)
bul(s, Inches(0.9), Inches(4.25), Inches(5.2), Inches(2.2), [
    "Time-domain: RMS, kurtosis, crest factor, peak-to-peak, skewness",
    "Freq-domain: FFT, spectral entropy, dominant freq, band energy",
    "PCA dimensionality reduction for cluster separability",
    "Clustering trained only on normal operating data",
    "Distance threshold-based anomaly scoring",
], 12, DARK)

shape(s, Inches(6.8), Inches(3.8), Inches(5.8), Inches(2.8), WHITE, GREEN, 0.02)
txt(s, Inches(7.1), Inches(3.9), Inches(5.2), Inches(0.3), "Results", 14, GREEN, True)
bul(s, Inches(7.1), Inches(4.25), Inches(5.2), Inches(2.2), [
    "Zero false positives on baseline validation",
    "WCSS compactness: 441.78 (tight clusters)",
    "Max distance from center: 5.86",
    "Average distance: 3.21",
    "Findings presented at defense R&D conference",
], 12, DARK, GREEN)


# ════════════════════════════════════════════════════
# SLIDE: CATTLE BIOMETRICS
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "Retrieval & Computer Vision \u2013 Goodhar", "Cattle Biometric Identification Platform")

# Problem
shape(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.9), WHITE, RGBColor(0xE7, 0x4C, 0x3C), 0.02)
txt(s, Inches(0.9), Inches(1.55), Inches(1.5), Inches(0.3), "Problem", 14, RGBColor(0xE7, 0x4C, 0x3C), True)
txt(s, Inches(2.5), Inches(1.55), Inches(9.8), Inches(0.7),
    "Traditional cattle ID (RFID, ear tags) is invasive, prone to loss/forgery, expensive at scale. "
    "Need non-invasive biometric identification across 100K+ cattle with poor-quality field images.",
    12, DARK)

# Pipeline flow
txt(s, Inches(0.8), Inches(2.6), Inches(5), Inches(0.3), "16-Stage Pipeline", 14, BLUE_DARK, True)
arrow_flow(s, Inches(0.8), Inches(2.95), Inches(11.5), [
    "Image\nInput", "YOLOv8\nDetection", "CLAHE\nEnhance", "Quality\nScoring", "Feature\nExtraction", "FAISS\nIndexing", "Retrieval\n& Ranking"
], [BLUE_DARK, BLUE, BLUE, BLUE_MED, BLUE_MED, BLUE_LIGHT, GREEN])

# Technical cards
shape(s, Inches(0.6), Inches(3.8), Inches(3.7), Inches(3.0), BG_CARD, BLUE_PALE, 0.02)
txt(s, Inches(0.9), Inches(3.9), Inches(3.2), Inches(0.3), "Preprocessing & Quality", 13, BLUE_DARK, True)
bul(s, Inches(0.9), Inches(4.2), Inches(3.2), Inches(2.4), [
    "rembg DL segmentation + GrabCut fallback",
    "CLAHE in LAB color space",
    "40+ metrics across 7 categories",
    "Confidence-based accept/reject",
], 11, DARK)

shape(s, Inches(4.6), Inches(3.8), Inches(3.7), Inches(3.0), BG_CARD, BLUE_PALE, 0.02)
txt(s, Inches(4.9), Inches(3.9), Inches(3.2), Inches(0.3), "Feature Extraction", 13, BLUE_DARK, True)
bul(s, Inches(4.9), Inches(4.2), Inches(3.2), Inches(2.4), [
    "ResNet, EfficientNet, ViT, ArcFace",
    "DINO, SimCLR, MAE, MoCo, BYOL",
    "512D L2-normalized embeddings",
    "20 parallel worker processes",
], 11, DARK)

shape(s, Inches(8.6), Inches(3.8), Inches(4.1), Inches(3.0), BG_CARD, BLUE_PALE, 0.02)
txt(s, Inches(8.9), Inches(3.9), Inches(3.5), Inches(0.3), "Vector Search & Fusion", 13, BLUE_DARK, True)
bul(s, Inches(8.9), Inches(4.2), Inches(3.5), Inches(2.4), [
    "FAISS multi-index per feature type",
    "6 fusion: intersection, weighted voting,",
    "  rank aggregation, cascade, union, adaptive",
    "Regional sharding (10\u201320x speedup)",
], 11, DARK)


# ════════════════════════════════════════════════════
# SLIDE: CATTLE - SCALABILITY
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "Cattle Biometric System \u2013 Scalability & Production Design")

stat(s, Inches(0.8), Inches(1.5), "100K+", "Images Processed", BLUE)
stat(s, Inches(3.5), Inches(1.5), "100K+", "FAISS Embeddings", BLUE_MED)
stat(s, Inches(6.2), Inches(1.5), "<3s", "End-to-End Latency", GREEN)
stat(s, Inches(8.9), Inches(1.5), "1M+", "Target Scale", BLUE_LIGHT)

# Architecture
shape(s, Inches(0.6), Inches(3.1), Inches(5.8), Inches(3.5), BG_CARD, BLUE_PALE, 0.02)
txt(s, Inches(0.9), Inches(3.2), Inches(5.2), Inches(0.3), "System Architecture", 14, BLUE_DARK, True)
bul(s, Inches(0.9), Inches(3.6), Inches(5.2), Inches(2.8), [
    "PostgreSQL for metadata + audit trail (soft deletes only)",
    "FAISS IndexIVFFlat with regional sharding",
    "S3/GCS for image storage (hot/cold/glacier)",
    "FastAPI + Docker for RESTful endpoints",
    "Monthly index rebuild with zero-downtime atomic swap",
    "Model versioning + embedding migration on retrain",
], 12, DARK)

shape(s, Inches(6.8), Inches(3.1), Inches(5.8), Inches(3.5), WHITE, GREEN, 0.02)
txt(s, Inches(7.1), Inches(3.2), Inches(5.2), Inches(0.3), "Decision Logic & Learning", 14, GREEN, True)
bul(s, Inches(7.1), Inches(3.6), Inches(5.2), Inches(2.8), [
    "High confidence (>0.88): automatic match",
    "Medium (0.80\u20130.88): manual review with side-by-side images",
    "Low (<0.80): flag as new cattle, prompt registration",
    "Active learning: operator feedback \u2192 hard example mining",
    "Monthly retraining: fine-tune SimCLR on new data",
    "Duplicate detection at registration (similarity >0.90)",
], 12, DARK, GREEN)


# ════════════════════════════════════════════════════
# SLIDE: RAG & AI AVATAR
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "RAG Systems & Real-Time AI Avatar")

# RAG
shape(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.5), WHITE, BLUE, 0.03)
txt(s, Inches(0.9), Inches(1.6), Inches(5.2), Inches(0.35), "RAG over Technical Report Corpora", 15, BLUE_DARK, True)
txt(s, Inches(0.9), Inches(2.0), Inches(1.2), Inches(0.25), "Problem:", 11, RGBColor(0xE7, 0x4C, 0x3C), True)
txt(s, Inches(2.1), Inches(2.0), Inches(4), Inches(0.25), "Engineers can't search multi-year internal reports efficiently", 11, DARK)
bul(s, Inches(0.9), Inches(2.4), Inches(5.2), Inches(1.3), [
    "Chunking + metadata filtering + embedding ingestion pipeline",
    "FAISS similarity search in FastAPI backend",
    "75% transcription latency reduction, 52% synthesis time reduction",
    "Control Gate: 100% responses scoped to local knowledge base",
], 11, DARK, GREEN)

# Avatar
shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5), WHITE, BLUE, 0.03)
txt(s, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.35), "Real-Time 3D AI Avatar (Defense)", 15, BLUE_DARK, True)
txt(s, Inches(7.1), Inches(2.0), Inches(1.2), Inches(0.25), "Problem:", 11, RGBColor(0xE7, 0x4C, 0x3C), True)
txt(s, Inches(8.3), Inches(2.0), Inches(4), Inches(0.25), "Need offline-capable conversational AI for defense deployments", 11, DARK)
bul(s, Inches(7.1), Inches(2.4), Inches(5.2), Inches(1.3), [
    "Whisper STT \u2192 Ollama LLM \u2192 LangChain RAG \u2192 Silero TTS",
    "STT: ~4s \u2192 <1s (75% via FP16), TTS: ~11.3s \u2192 ~5.4s (52%)",
    "Three.js 3D with smooth viseme blending + facial animation",
    "FastAPI + React + WebSocket, fully offline-capable",
], 11, DARK, GREEN)

# Avatar pipeline flow
txt(s, Inches(0.8), Inches(4.2), Inches(5), Inches(0.3), "AI Avatar System Flow", 14, BLUE_DARK, True)
arrow_flow(s, Inches(0.8), Inches(4.6), Inches(11.5), [
    "User\nSpeech", "Whisper\nSTT (FP16)", "Ollama\nLLM", "LangChain\nRAG Gate", "Silero\nTTS", "Three.js\nLip-Sync"
], [BLUE_DARK, BLUE, BLUE_MED, BLUE_LIGHT, BLUE_MED, GREEN])

# Agritech Advisory
shape(s, Inches(0.6), Inches(5.5), Inches(12), Inches(1.5), BG_CARD, BLUE_PALE, 0.02)
txt(s, Inches(0.9), Inches(5.6), Inches(10), Inches(0.35), "AI/ML Advisory \u2013 Agritech Startup  |  Mar 2025 \u2013 Aug 2025", 14, BLUE_DARK, True)
bul(s, Inches(0.9), Inches(5.95), Inches(11.2), Inches(0.9), [
    "Benchmarked OpenAI, Anthropic, LLaMA, Mistral for agriculture QA  \u2022  Bilingual English/Telugu for Andhra Pradesh farmers",
    "Designed STT fine-tuning pipeline: audio collection \u2192 transcription \u2192 diarization \u2192 agriculture-specific vocabulary tuning",
], 12, DARK)


# ════════════════════════════════════════════════════
# SECTION: EDUCATION
# ════════════════════════════════════════════════════
section("Education & Credentials")

s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "Education & Certifications")

shape(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.8), WHITE, BLUE, 0.03)
txt(s, Inches(0.9), Inches(1.6), Inches(5.2), Inches(0.35), "B.Tech \u2013 Mechanical Engineering", 18, BLUE_DARK, True)
txt(s, Inches(0.9), Inches(2.0), Inches(5.2), Inches(0.25), "Mahindra University, Hyderabad  |  2018 \u2013 2022  |  CGPA: 7.5", 13, BLUE_MED)
txt(s, Inches(0.9), Inches(2.4), Inches(5.2), Inches(1.5),
    "Foundation in computational modeling, numerical methods (Euler-Maruyama, GFBM, Monte Carlo), "
    "systems analysis, and GPU programming. This engineering background is what enables my approach to ML \u2014 "
    "I think in systems, not just models.",
    12, DARK)

shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.8), WHITE, BLUE, 0.03)
txt(s, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.35), "Certifications", 18, BLUE_DARK, True)
bul(s, Inches(7.1), Inches(2.1), Inches(5.2), Inches(2.0), [
    "Stanford Machine Learning (Andrew Ng)",
    "Deep Learning Specialization (deeplearning.ai)",
    "IBM Machine Learning Professional Certificate",
    "Applied ML (University of Michigan)",
    "Deep Learning with TensorFlow",
], 13, DARK)

shape(s, Inches(0.6), Inches(4.6), Inches(12), Inches(1.2), BG_CARD, BLUE_PALE, 0.02)
txt(s, Inches(0.9), Inches(4.7), Inches(10), Inches(0.35), "Relevant Coursework", 15, BLUE_DARK, True)
txt(s, Inches(0.9), Inches(5.1), Inches(11), Inches(0.55),
    "Intro to Computer Science  \u2022  Linear Algebra & Matrices  \u2022  Probability & Statistics  \u2022  Data Structures  \u2022  "
    "Big Data Computing  \u2022  Advanced Data Analytics  \u2022  Time Series Forecasting  \u2022  GPU Programming  \u2022  Computer Aided Engineering Design",
    12, DARK)

# Academic projects
txt(s, Inches(0.8), Inches(6.1), Inches(5), Inches(0.3), "Academic Research", 14, BLUE_DARK, True)
txt(s, Inches(0.8), Inches(6.4), Inches(12), Inches(0.5),
    "CFD Aircraft Simulation (DLR F11, 20M-element mesh, NASA-validated)  \u2022  Stock Forecasting (Euler-Maruyama, Monte Carlo)  \u2022  "
    "GFBM with Hurst Exponent  \u2022  Quantum Neutrino Oscillations  \u2022  Movie Recommender (Hadoop MapReduce)  \u2022  Wind Turbine Prediction (Flask, RF >92% R\u00b2)",
    11, GRAY)


# ════════════════════════════════════════════════════
# SLIDE: CORE SKILLS
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "Core Skills & Technologies")

skills = [
    ("Machine Learning & AI", "LLMs, RAG, Deep Learning, Feature Engineering, Model Optimization, Anomaly Detection, NLP, Computer Vision"),
    ("Retrieval & Vector Systems", "FAISS, Embedding Pipelines, Vector Search, ChromaDB, Multi-Index Fusion, Ranking Strategies"),
    ("Cloud & Infrastructure", "AWS (EC2, Lambda, SQS, IAM, ECS, Amplify, GuardDuty, Secrets Mgr), Terraform, Docker, CI/CD, SOC2"),
    ("Data & Backend", "PySpark, ETL, FastAPI, SQL, PostgreSQL, Hadoop, MapReduce, Async Processing, WebSockets"),
    ("Simulation & Engineering", "ANSYS APDL, MOSES, OrcaFlex, FreeCAD, OpenCV, Signal Processing, CFD (ANSYS Fluent)"),
    ("Programming", "Python (primary), C, Go, Java, Flask, Django, React, Three.js"),
]

for i, (cat, items) in enumerate(skills):
    row = i // 2; col = i % 2
    x = Inches(0.6) + Inches(col * 6.2)
    y = Inches(1.5) + Inches(row * 1.75)
    shape(s, x, y, Inches(5.8), Inches(1.55), WHITE, BLUE, 0.03)
    txt(s, x + Inches(0.3), y + Inches(0.1), Inches(5.2), Inches(0.35), cat, 15, BLUE_DARK, True)
    txt(s, x + Inches(0.3), y + Inches(0.5), Inches(5.2), Inches(0.85), items, 12, DARK)


# ════════════════════════════════════════════════════
# SLIDE: IMPACT
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "Impact at a Glance")

impacts = [
    ("88,560", "Simulated sea states", "Digital Twin"),
    ("~1.08M", "Model fits executed", "ML Experimentation"),
    ("100K+", "Cattle images processed", "Goodhar Platform"),
    ("30%", "AWS cost reduction", "EC2 \u2192 Lambda"),
    ("5\u20137x", "Batch speed improvement", "Lambda migration"),
    ("41%", "Angular error reduction", "Circular regression"),
    ("75%", "STT latency reduction", "AI Avatar (FP16)"),
    ("52%", "TTS latency reduction", "Phonetic recognizer"),
    ("97.17%", "Direction prediction", "Random Forest"),
    (">99%", "Migration parity", "SAS \u2192 PySpark"),
    ("0", "False positives", "Anomaly Detection"),
    ("80%+", "Manual work reduced", "Simulation Automation"),
]

for i, (val, label, ctx) in enumerate(impacts):
    row = i // 4; col = i % 4
    x = Inches(0.5) + Inches(col * 3.15)
    y = Inches(1.5) + Inches(row * 1.8)
    shape(s, x, y, Inches(2.9), Inches(1.55), WHITE, BLUE, 0.04)
    txt(s, x + Inches(0.1), y + Inches(0.1), Inches(2.7), Inches(0.55), val, 28, BLUE, True, PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), y + Inches(0.65), Inches(2.7), Inches(0.35), label, 12, DARK, False, PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), y + Inches(1.05), Inches(2.7), Inches(0.3), ctx, 10, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE: WHAT I BRING
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
header(s, "What I Bring to the Table")

strengths = [
    ("End-to-End Systems Builder", "I build complete systems, not just models. From raw data ingestion to SOC2-compliant production deployment \u2014 infrastructure, pipelines, monitoring, and all."),
    ("Cross-Domain Problem Solver", "Defense, offshore, agritech, patent law, blockchain. Each domain taught me to adapt fast, work with domain experts, and ship under real constraints."),
    ("Scale-Tested Engineering", "88K+ simulations, 1M+ model fits, 100K+ image embeddings. I design for production loads from day one, not as an afterthought."),
    ("Research \u2192 Production Bridge", "Mechanical engineering foundation + ML expertise = I think in systems. I take research prototypes and turn them into maintainable, deployed infrastructure."),
]

for i, (title, desc) in enumerate(strengths):
    y = Inches(1.5) + Inches(i * 1.35)
    shape(s, Inches(0.6), y, Inches(11.8), Inches(1.15), WHITE, BLUE, 0.03)
    txt(s, Inches(1.0), y + Inches(0.1), Inches(3.2), Inches(0.35), title, 16, BLUE_DARK, True)
    txt(s, Inches(1.0), y + Inches(0.5), Inches(11), Inches(0.55), desc, 13, DARK)


# ════════════════════════════════════════════════════
# SLIDE: CONTACT
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
rect(s, Inches(0), Inches(5.5), Inches(13.333), Inches(2), BLUE_DARK)

try:
    s.shapes.add_picture(PHOTO, Inches(5.6), Inches(0.8), Inches(2.1), Inches(2.1))
except:
    pass

txt(s, Inches(1), Inches(3.1), Inches(11.3), Inches(0.7), "Let's Build Something Together", 40, BLUE_DARK, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8),
    "Open to full-time roles, contract engagements, consulting opportunities, and\n"
    "partnering with founders building at the intersection of AI and real-world systems.",
    15, GRAY, False, PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(5.7), Inches(11.3), Inches(0.5), "sriharisirisipalli0@gmail.com", 20, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.2), Inches(11.3), Inches(0.4), "srihari-sirisipalli.github.io  |  calendly.com/sriharisirisipalli0", 14, BLUE_LIGHT, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.4), "LinkedIn: sri-hari-sirisipalli  |  GitHub: srihari-sirisipalli", 13, RGBColor(0x90, 0xA8, 0xC0), False, PP_ALIGN.CENTER)


# ── Save ──
output = r"C:\Users\Sri\Documents\My Documents\Sri_Hari_Pitch_Deck.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
